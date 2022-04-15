from pptx import Presentation
from pptx.util import Cm

#Presentaionオブジェクトを生成
ppt = Presentation()

#スライドのサイズを指定
ppt.slide_width = Cm(33.867)
ppt.slide_height = Cm(19.05)

#追加するスライドを選択
slide_layout_0 = ppt.slide_layouts[0]

#追加するスライドを選択
slide_layout_0 = ppt.slide_layouts[0]
#スライドを追加#
slide_0 = ppt.slides.add_slide(slide_layout_0)
slide_0_title = slide_0.placeholders[0]
print(type(slide_0_title))
slide_0_title.width = Cm(25.4)
slide_0_title.height = Cm(6.63)
slide_0_title.text = "HOGEHOGE"

slide_0_subtitle = slide_0.placeholders[1]
slide_0_subtitle.text = "FOOFOO"
slide_0_subtitle.text += "FOOFOO"

ppt.save("./test.pptx")